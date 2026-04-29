from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("presentation")
PPTX_PATH = OUT_DIR / "glitch_museum_mode_peer_presentation.pptx"
NOTES_PATH = OUT_DIR / "glitch_museum_mode_peer_speaker_notes.md"

NAVY = RGBColor(18, 31, 46)
INK = RGBColor(30, 41, 59)
MUTED = RGBColor(91, 105, 125)
TEAL = RGBColor(20, 184, 166)
AMBER = RGBColor(245, 158, 11)
WHITE = RGBColor(255, 255, 255)
SOFT = RGBColor(241, 245, 249)


def add_textbox(slide, left, top, width, height, text, size=28, color=INK, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = "Aptos"
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.7), Inches(0.48), Inches(11.8), Inches(0.72), title, 34, INK, True)
    if subtitle:
        add_textbox(slide, Inches(0.74), Inches(1.16), Inches(11.1), Inches(0.38), subtitle, 16, MUTED)


def add_tag(slide, left, top, width, text, fill_color=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Aptos"


def add_bullets(slide, items, left, top, width, height, size=23):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.size = Pt(size)
        paragraph.font.name = "Aptos"
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(11)


def add_footer(slide, text):
    add_textbox(slide, Inches(0.75), Inches(6.92), Inches(11.8), Inches(0.25), text, 10, MUTED)


def build_deck():
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    notes = []

    # Slide 1
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_tag(slide, Inches(0.8), Inches(0.75), Inches(2.35), "3 minute share", TEAL)
    add_textbox(slide, Inches(0.78), Inches(1.65), Inches(11.5), Inches(1.05), "Glitch Museum Mode", 54, WHITE, True)
    add_textbox(slide, Inches(0.82), Inches(2.78), Inches(10.4), Inches(0.85), "Turning old AI-generated bugs into evidence-backed exhibits.", 27, TEAL, True)
    add_textbox(slide, Inches(0.85), Inches(4.0), Inches(9.4), Inches(0.76), "A Streamlit guessing game plus a local RAG guide that explains what broke, how it was fixed, and why the tests matter.", 22, SOFT)
    add_footer(slide, "Peer presentation focus: favorite parts, not every grading requirement.")
    notes.append(
        "Open with the project name. Say: I turned a fixed AI-generated guessing game into a RAG museum where old bugs become exhibits. The Loom is my full grading proof; this short presentation is about the parts I found most interesting."
    )

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Favorite Part: Bugs Became Exhibits", "The RAG guide explains the project's debugging history with evidence.")
    add_tag(slide, Inches(0.85), Inches(1.95), Inches(2.75), "Demo this", TEAL)
    add_textbox(slide, Inches(0.9), Inches(2.65), Inches(5.2), Inches(1.15), "Shapeshifting Secret Number", 34, NAVY, True)
    add_bullets(
        slide,
        [
            "Select the artifact in Glitch Museum Mode.",
            "Show the answer generated from retrieved project files.",
            "Point out confidence and evidence from reflection/code.",
        ],
        Inches(6.35),
        Inches(2.05),
        Inches(5.85),
        Inches(3.4),
        22,
    )
    add_textbox(slide, Inches(0.92), Inches(4.2), Inches(5.0), Inches(1.0), "Why it matters: the answer is grounded in this project, not generic AI advice.", 24, INK, True)
    add_footer(slide, "Suggested live demo: Inspect Artifact -> Shapeshifting Secret Number")
    notes.append(
        "Show the Shapeshifting Secret Number artifact. Say this is my favorite part because it turns a bug into an exhibit. The system retrieves evidence from project files and explains why Streamlit reruns made the secret number unstable."
    )

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Most Responsible Part: It Can Say 'I Don't Know'", "The guardrail is more interesting than a perfect-looking answer.")
    add_textbox(slide, Inches(0.9), Inches(1.92), Inches(4.7), Inches(1.0), "Unrelated question:", 26, MUTED, True)
    add_textbox(slide, Inches(0.9), Inches(2.72), Inches(5.7), Inches(0.95), "Explain database migrations and cloud billing", 27, NAVY, True)
    add_bullets(
        slide,
        [
            "No matching project evidence",
            "Confidence drops to 0.00",
            "The guide returns a fallback instead of pretending",
        ],
        Inches(7.05),
        Inches(2.0),
        Inches(5.25),
        Inches(3.2),
        23,
    )
    add_footer(slide, "Responsible AI moment: low confidence is a feature, not a failure.")
    notes.append(
        "Demo or describe the unrelated question. Say: A weaker AI system might answer anyway. Mine refuses when it cannot retrieve project evidence. That made me think about confidence and fallback behavior as part of responsible AI."
    )

    # Slide 4
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SOFT
    add_title(slide, "What I Learned", "AI engineering is verification, not just generation.")
    add_bullets(
        slide,
        [
            "AI-generated code can run and still be logically wrong.",
            "RAG is valuable when evidence shapes the answer.",
            "Tests and confidence scoring make demos more trustworthy.",
            "Human review still matters when AI suggestions are incomplete.",
        ],
        Inches(0.9),
        Inches(1.85),
        Inches(7.25),
        Inches(4.4),
        24,
    )
    add_textbox(slide, Inches(8.55), Inches(2.1), Inches(3.7), Inches(2.4), "My takeaway:\nBuild AI that is useful, testable, and honest about its limits.", 25, NAVY, True)
    add_footer(slide, "Leave about one minute for questions.")
    notes.append(
        "Close with the lesson: AI code can look right but behave wrong. This project taught me to pair AI features with evidence, tests, confidence, and human review."
    )

    prs.save(PPTX_PATH)
    NOTES_PATH.write_text(format_notes(notes), encoding="utf-8")
    print(PPTX_PATH)
    print(NOTES_PATH)


def format_notes(notes):
    lines = ["# Peer Presentation Speaker Notes", ""]
    for i, note in enumerate(notes, start=1):
        lines.extend([f"## Slide {i}", "", note, ""])
    return "\n".join(lines)


if __name__ == "__main__":
    build_deck()
